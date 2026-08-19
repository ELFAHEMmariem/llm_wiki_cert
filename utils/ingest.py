import base64
import io
import os
import re
import sys
import tempfile
import zipfile
import tarfile
from pathlib import Path
from typing import Optional

from bs4 import BeautifulSoup
import docx
import fitz  # PyMuPDF
import requests
import torch
import whisper
import yt_dlp
from dotenv import load_dotenv
from groq import Groq
from pptx import Presentation

# Force le rechargement des variables d'environnement
load_dotenv(override=True)

# Correction d'encodage pour la console Windows
if sys.platform.startswith("win"):
    sys.stdout = io.TextIOWrapper(
        sys.stdout.buffer, encoding="utf-8", errors="replace"
    )
    sys.stderr = io.TextIOWrapper(
        sys.stderr.buffer, encoding="utf-8", errors="replace"
    )


class IngestPipeline:

    def __init__(
        self,
        api_key: Optional[str] = None,
        output_dir: str = "wiki",
        llm_model: str = "openai/gpt-oss-120b",
    ):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self.api_key = api_key or os.getenv("GROQ_API_KEY")
        if self.api_key:
            self.client = Groq(api_key=self.api_key)
        else:
            self.client = None

        self.llm_model = llm_model
        self.whisper_model = None
        has_cuda = torch.cuda.is_available()
        self.device = "cuda" if has_cuda else "cpu"

    # =========================================================
    # APPEL GÉNÉRIQUE CHAT VIA GROQ (SANS MODÈLES OBSOLÈTES)
    # =========================================================
    def _call_groq_chat(self, prompt: str, model: Optional[str] = None) -> str:
        """Appel générique à l'API Groq avec cascade de repli (fallback)."""
        if not self.client:
            raise ValueError(
                "La clé GROQ_API_KEY est introuvable dans le fichier .env"
            )

        target_model = model or self.llm_model
        
        # Modèles valides et actuellement supportés sur l'API Groq
        models_to_try = [
            target_model,
            "llama-3.1-70b-versatile",
            "llama3-70b-8192",
            "llama3-8b-8192",
            "mixtral-8x7b-32768",
        ]

        # Déduplication tout en conservant l'ordre
        seen = set()
        models_to_try = [m for m in models_to_try if not (m in seen or seen.add(m))]

        last_error = None
        for m in models_to_try:
            try:
                response = self.client.chat.completions.create(
                    model=m,
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.2,
                    max_tokens=2048,
                )
                return response.choices[0].message.content
            except Exception as e:
                last_error = e
                print(f"[AVERTISSEMENT GROQ] Échec avec le modèle {m} : {e}")

        raise RuntimeError(
            f"Aucun modèle Groq n'a pu répondre. Dernier détail : {last_error}"
        )

    # =========================================================
    # ROUTEUR UNIVERSEL DE FORMATS DE FICHIERS
    # =========================================================
    def process_any_file(
        self, file_path: Path, context_tag: str = "Ingestion"
    ) -> str:
        """Route n'importe quel fichier vers l'extracteur approprié."""
        ext = file_path.suffix.lower()
        file_str = str(file_path)

        # 1. Documents PDF
        if ext == ".pdf":
            return self.process_pdf_smart(file_str, context_tag)

        # 2. Images (Vision)
        elif ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".gif"]:
            return self.process_image(file_str, context_tag)

        # 3. Documents Bureautique
        elif ext in [".docx", ".doc"]:
            return self.process_docx(file_str, context_tag)
        elif ext in [".pptx", ".ppt"]:
            return self.process_pptx(file_str, context_tag)

        # 4. Audio et Vidéo
        elif ext in [
            ".mp3",
            ".wav",
            ".mp4",
            ".m4a",
            ".aac",
            ".flac",
            ".ogg",
            ".mkv",
            ".avi",
            ".mov",
        ]:
            return self.process_audio_video(file_str, context_tag)

        # 5. Raccourcis et Liens Web
        elif ext == ".url":
            return self.process_url_file(file_path, context_tag)

        # 6. Archives (.zip, .tar, .gz)
        elif ext in [".zip", ".tar", ".gz", ".tgz"]:
            return self.process_archive(file_path, context_tag)

        # 7. Fichiers texte brut / Code / Structurés
        elif ext in [
            ".txt",
            ".md",
            ".csv",
            ".json",
            ".xml",
            ".html",
            ".htm",
            ".py",
            ".js",
            ".c",
            ".cpp",
            ".java",
            ".log",
            ".yaml",
            ".yml",
        ]:
            return file_path.read_text(encoding="utf-8", errors="ignore")

        # 8. FALLBACK UNIVERSEL
        else:
            print(
                f"[FALLBACK UNIVERSEL] Format non standard ({ext}) -> Tentative"
                " d'extraction brute..."
            )
            try:
                content = file_path.read_text(encoding="utf-8", errors="ignore")
                if len(content.strip()) > 20:
                    return content
            except Exception:
                pass
            return (
                "[Format non lisible directement] Nom du fichier :"
                f" {file_path.name}"
            )

    # =========================================================
    # EXTRACTION INTELLIGENTE DE PDF
    # =========================================================
    def process_pdf_smart(self, file_path: str, context_tag: str) -> str:
        print(f"[PDF-SMART] Extraction PDF : {os.path.basename(file_path)}...")
        full_text = []
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text().strip()
                if text:
                    full_text.append(f"--- Page {page_num + 1} ---\n{text}")
            doc.close()
            return "\n\n".join(full_text)
        except Exception as e:
            print(f"[ERREUR] Échec extraction PDF : {e}")
            return f"Erreur lors de la lecture du PDF : {e}"

    # =========================================================
    # TRAITEMENT DES IMAGES VIA GROQ VISION
    # =========================================================
    def process_image(self, file_path: str, context_tag: str) -> str:
        print(
            "[IMAGE-OCR] Analyse image"
            f" ({context_tag}) : {os.path.basename(file_path)}..."
        )
        if not self.client:
            return "Client Groq indisponible pour la vision."
        try:
            with open(file_path, "rb") as image_file:
                encoded_string = base64.b64encode(image_file.read()).decode(
                    "utf-8"
                )

            ext = Path(file_path).suffix.lower().replace(".", "")
            mime_type = "image/png" if ext == "png" else "image/jpeg"

            response = self.client.chat.completions.create(
                model="openai/gpt-oss-120b",
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "Transcris fidèlement tout le texte présent"
                                " dans cette image. Décris également les"
                                " schémas, diagrammes ou graphiques s'il y en"
                                " a."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": (
                                    f"data:{mime_type};base64,{encoded_string}"
                                )
                            },
                        },
                    ],
                }],
                temperature=0.2,
            )
            return response.choices[0].message.content

        except Exception as e:
            print(f"[ERREUR] OCR Image Groq : {e}")
            # Fallback en mode texte si l'API vision échoue
            return f"Erreur lors de l'analyse d'image : {e}"

    # =========================================================
    # TRAITEMENT DES ARCHIVES (.ZIP, .TAR, .GZ)
    # =========================================================
    def process_archive(self, file_path: Path, context_tag: str) -> str:
        print(f"[ARCHIVE] Extraction de l'archive : {file_path.name}")
        extracted_texts = []
        temp_dir = tempfile.mkdtemp()

        try:
            if file_path.suffix.lower() == ".zip":
                with zipfile.ZipFile(file_path, "r") as zip_ref:
                    zip_ref.extractall(temp_dir)
            elif file_path.suffix.lower() in [".tar", ".gz", ".tgz"]:
                with tarfile.open(file_path, "r:*") as tar_ref:
                    tar_ref.extractall(temp_dir)

            for sub_file in Path(temp_dir).rglob("*"):
                if sub_file.is_file():
                    text = self.process_any_file(sub_file, context_tag)
                    extracted_texts.append(
                        f"--- Fichier inclus : {sub_file.name} ---\n{text}"
                    )

        except Exception as e:
            print(f"[ERREUR] Extraction archive : {e}")
            return f"Erreur lors de la décompression de l'archive : {e}"

        return "\n\n".join(extracted_texts)

    # =========================================================
    # TRAITEMENT DES URL & FICHIERS .URL
    # =========================================================
    def process_url_file(self, file_path: Path, context_tag: str) -> str:
        try:
            content = file_path.read_text(encoding="utf-8", errors="ignore")
            url_match = re.search(r"URL=(.*)", content, re.IGNORECASE)
            target_url = (
                url_match.group(1).strip() if url_match else content.strip()
            )

            if target_url.startswith("http"):
                return self.process_url_smart(target_url, context_tag)
            return f"Lien URL invalide : {target_url}"
        except Exception as e:
            return f"Erreur de lecture du fichier URL : {e}"

    def process_url_smart(self, url: str, context_tag: str) -> str:
        print(f"[URL-INGEST] Analyse de l'URL : {url}")

        media_domains = [
            "youtube.com",
            "youtu.be",
            "vimeo.com",
            "soundcloud.com",
            "dailymotion.com",
        ]
        is_media_platform = any(
            domain in url.lower() for domain in media_domains
        )

        if is_media_platform:
            return self.process_url_media(url, context_tag)

        try:
            headers = {
                "User-Agent": (
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/115.0.0.0 Safari/537.36"
                )
            }
            response = requests.get(url, headers=headers, timeout=12)
            response.raise_for_status()

            content_type = response.headers.get("Content-Type", "").lower()

            if (
                "text/html" in content_type
                or "application/xhtml+xml" in content_type
            ):
                soup = BeautifulSoup(response.text, "html.parser")

                for element in soup([
                    "script",
                    "style",
                    "nav",
                    "footer",
                    "header",
                    "aside",
                    "noscript",
                ]):
                    element.decompose()

                lines = (line.strip() for line in soup.get_text().splitlines())
                chunks = (
                    phrase.strip()
                    for line in lines
                    for phrase in line.split("  ")
                )
                extracted_text = "\n".join(chunk for chunk in chunks if chunk)

                if len(extracted_text.strip()) > 100:
                    print(
                        "[URL-TEXT] Texte extrait avec succès depuis la page"
                        " Web."
                    )
                    title = (
                        soup.title.string.strip() if soup.title else "Page Web"
                    )
                    return (
                        f"--- CONTENU DE LA PAGE WEB : {title} ({url}) ---\n\n{extracted_text}"
                    )

        except Exception as e:
            print(
                f"[URL-TEXT] Échec de la récupération texte HTML ({e})."
                " Tentative de fallback Média..."
            )

        return self.process_url_media(url, context_tag)

    def process_url_media(self, url: str, context_tag: str) -> str:
        print(
            "[URL-MEDIA] Téléchargement média via yt-dlp :"
            f" {url} ..."
        )
        temp_dir = tempfile.mkdtemp()
        output_template = os.path.join(temp_dir, "downloaded_audio.%(ext)s")

        ydl_opts = {
            "format": "bestaudio/best",
            "outtmpl": output_template,
            "postprocessors": [{
                "key": "FFmpegExtractAudio",
                "preferredcodec": "mp3",
                "preferredquality": "192",
            }],
            "quiet": True,
            "no_warnings": True,
        }

        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)
                title = info.get("title", "Titre inconnu")

            audio_file = os.path.join(temp_dir, "downloaded_audio.mp3")
            if not os.path.exists(audio_file):
                raise FileNotFoundError("Échec de la conversion de l'audio.")

            print(f"[WHISPER] Transcription de '{title}'...")
            transcription = self.process_audio_video(audio_file, context_tag)

            if os.path.exists(audio_file):
                os.remove(audio_file)
            os.rmdir(temp_dir)

            return (
                f"--- TRANSCRIPTION MÉDIA : {title} ({url}) ---\n\n{transcription}"
            )

        except Exception as e:
            print(f"[ERREUR] Échec du traitement URL Média {url} : {e}")
            return f"Erreur lors du traitement de l'URL ({url}) : {e}"

    # =========================================================
    # EXTRACTION DOCX, PPTX & AUDIO LOCAL
    # =========================================================
    def process_docx(self, file_path: str, context_tag: str) -> str:
        try:
            doc = docx.Document(file_path)
            full_content = []

            for p in doc.paragraphs:
                text = p.text.strip()
                if text:
                    full_content.append(text)

            for table in doc.tables:
                table_md = []
                for i, row in enumerate(table.rows):
                    row_cells = [
                        cell.text.strip().replace("\n", " ")
                        for cell in row.cells
                    ]
                    table_md.append("| " + " | ".join(row_cells) + " |")
                    if i == 0:
                        table_md.append(
                            "| " + " | ".join(["---"] * len(row_cells)) + " |"
                        )
                full_content.append("\n".join(table_md))

            return "\n\n".join(full_content)
        except Exception as e:
            return f"Erreur d'extraction DOCX : {e}"

    def process_pptx(self, file_path: str, context_tag: str) -> str:
        try:
            prs = Presentation(file_path)
            slides_text = []

            for idx, slide in enumerate(prs.slides, 1):
                slide_lines = [f"Diapositive {idx}:"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_lines.append(f"- {text}")
                slides_text.append("\n".join(slide_lines))

            return "\n\n".join(slides_text)
        except Exception as e:
            return f"Erreur d'extraction PPTX : {e}"

    def process_audio_video(self, file_path: str, context_tag: str) -> str:
        try:
            if self.whisper_model is None:
                self.whisper_model = whisper.load_model(
                    "base", device=self.device
                )

            result = self.whisper_model.transcribe(file_path)
            raw_text = result["text"].strip()
            paragraphs = re.split(r"(?<=[.!?]) +", raw_text)

            formatted_text = ""
            current_p = []

            for sentence in paragraphs:
                current_p.append(sentence)
                if len(" ".join(current_p)) > 300:
                    formatted_text += " ".join(current_p) + "\n\n"
                    current_p = []

            if current_p:
                formatted_text += " ".join(current_p)

            return formatted_text
        except Exception as e:
            return f"Erreur de transcription Audio/Vidéo : {e}"

    # =========================================================
    # MÉTHODE PRINCIPALE D'INGESTION PAR FICHIER
    # =========================================================
    def process_file(self, file_path: Path) -> Optional[Path]:
        """Prétraitement universel et génération de la fiche Wiki Markdown."""
        if not file_path.exists():
            return None

        print(f"⚙️ [INGESTION] Traitement de : {file_path.name}")
        raw_text = self.process_any_file(file_path, context_tag=file_path.stem)

        prompt_structuration = f"""Tu es un analyste de sécurité informatique.
Transforme l'extraction brute ci-dessous en une fiche Wiki structurée et claire en Markdown.

NOM DU FICHIER : {file_path.name}

CONTENU BRUT EXTRAIT :
{raw_text[:20000]}

Consignes :
1. Génère un titre principal `# Fiche : {file_path.stem}`.
2. Identifie clairement les métadonnées clés (Cadre/Événement, Présentateur/Intervenants/Organismes).
3. Rédige un résumé synthétique des informations clés.
4. Si des listes d'invités/intervenants sont présentes, rassemble-les dans un tableau Markdown synthétique.
"""

        try:
            fiche_md = self._call_groq_chat(prompt_structuration)
            output_file = self.output_dir / f"{file_path.stem}.md"
            output_file.write_text(fiche_md, encoding="utf-8")
            print(f"✅ Fiche générée avec succès : {output_file}")
            return output_file
        except Exception as e:
            print(f"❌ Erreur lors de la génération de la fiche Markdown : {e}")
            return None